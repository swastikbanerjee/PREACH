from intel_extension_for_transformers.neural_chat import build_chatbot
from sklearn.feature_extraction.text import TfidfVectorizer
import emojis
import nltk
import re
import string
import google.generativeai as genai
from nltk.corpus import stopwords
from gensim import corpora, models
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from gensim.models import CoherenceModel
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import transformers
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os
import moviepy.editor as mp
import speech_recognition as sr
import transformers
import warnings
warnings.filterwarnings("ignore")
import streamlit as st
import tempfile
# nltk.download('stopwords')
# nltk.download('punkt')
# nltk.download('wordnet')

model_name = 'Intel/neural-chat-7b-v3-1'
GOOGLE_API_KEY = 'AIzaSyBileOLS4Ys9Nk1X27OqBccsEmgxOWmV54'
model = transformers.AutoModelForCausalLM.from_pretrained(model_name)
tokenizer = transformers.AutoTokenizer.from_pretrained(model_name)
genai.configure(api_key=GOOGLE_API_KEY)
modelG = genai.GenerativeModel('gemini-pro')
chatbot = build_chatbot()

# Function to convert audio to transcript using SpeechRecognition
def convert_audio_to_transcript(audio_path):
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio_data = recognizer.record(source)
        transcript = recognizer.recognize_google(audio_data)
    return transcript
# Function to convert video to transcript
def convert_video_to_transcript(video_path):
    # Convert video to audio
    audio_path = 'temp_audio.wav'
    video = mp.VideoFileClip(video_path)
    audio = video.audio
    audio.write_audiofile(audio_path)
    # Convert audio to transcript
    transcript = convert_audio_to_transcript(audio_path)
    # Delete temporary audio file
    os.remove(audio_path)
    return transcript
# Main function to handle different types of input
def handle_input(input_data):
    # Detect input type
    ext = input_data.name.split(".")[-1]
    temp_dir = tempfile.mkdtemp()
    # Save the uploaded file to the temporary directory
    file_path = os.path.join(temp_dir, input_data.name)
    with open(file_path, "wb") as f:
        f.write(input_data.read())
        
    if ext == 'mp4' or ext == 'avi':
        # Input is a video file
        transcript = convert_video_to_transcript(file_path)
    elif ext == 'wav' or ext == 'mp3':
        # Input is an audio file
        transcript = convert_audio_to_transcript(file_path)
    elif ext == 'txt':
        # Input is a text file
        with open(file_path, 'r') as file:
            transcript = file.read()
    else:
        # Input is text
        transcript = input_data
    return transcript
def preprocess_textS(text): #for summary
    # Remove emojis
    text = emojis.decode(text)
    # Tokenize the text into words
    words = nltk.word_tokenize(text)
    # Join the words back into a single string
    preprocessed_text = " ".join(words)
    return preprocessed_text
def preprocess_textK(text):#for keyword generation
    # Remove emojis
    text = emojis.decode(text)
    # Tokenize the text into words
    words = nltk.word_tokenize(text)
    # Remove punctuation, stopwords, and non-alphabetic words
    stop_words = set(stopwords.words('english'))
    words = [word for word in words if word.lower() not in stop_words and word not in string.punctuation and word.isalpha()]
    # Lemmatize words
    lemmatizer = WordNetLemmatizer()
    words = [lemmatizer.lemmatize(word) for word in words]
    return words
def preprocess_textW(text): #for wordcloud
    # Remove emojis
    text = emojis.decode(text)
    # Tokenize the text into words
    words = nltk.word_tokenize(text)
    # Remove punctuation
    words = [word for word in words if word not in string.punctuation]
    # Remove stopwords
    stop_words = set(stopwords.words('english'))
    words = [word for word in words if word.lower() not in stop_words]
    # Convert words to lowercase
    words = [word.lower() for word in words]
    # Join the words back into a single string
    preprocessed_text = " ".join(words)
    return preprocessed_text
def calculate_percentage_of_words(transcript_text):
    # Split the transcript text into words
    words = transcript_text.split()
    # Count the total number of words
    total_words = len(words)
    percentage=15
    # Calculate 15% of the total words
    fifteen_percent = int(total_words * (percentage / 100))
    return fifteen_percent
def generate_response(system_input, user_input): #neural chat
    # Format the input using the provided template
    prompt = f"### System:\n{system_input}\n### User:\n{user_input}\n### Assistant:\n"
    # Tokenize and encode the prompt
    inputs = tokenizer.encode(prompt, return_tensors="pt", add_special_tokens=False)
    # Generate a response
    outputs = model.generate(inputs, max_length=3500, num_return_sequences=1, pad_token_id = tokenizer.eos_token_id)
    response = tokenizer.decode(outputs[0], skip_special_tokens=True)
    # Extract only the assistant's response
    return response.split("### Assistant:\n")[-1]
def generate_summary(text): #Summary for ppt
    text=preprocess_textS(text)
#     question = f"Can you generate a 50 word summary for the following paragraph: {text}?"
#     response = chatbot.predict(question)
    response = modelG.generate_content(f"Can you generate a 50 word summary for the following paragraph: {text}?")
    return response.text
def generate_summary_ext(text,fifteen_percent): #Summary for ppt 15%
    text=preprocess_textS(text)
#     question = f"Can you generate a {fifteen_percent} word summary for the following paragraph: {text}?"
#     response = chatbot.predict(question)
    response = modelG.generate_content(f"Can you generate a {fifteen_percent} word summary for the following paragraph: {text}?")
    return response.text
def generate_title(text): #title
#     system_input = "You are a creative writing assistant. Your mission is to help users generate beautiful and thought provoking powerpoint presentation titles based on the text they input. Generate a title with a very short description of its meaning below. Format like-Title: ... Meaning: ..."
#     question = f"Can you generate a title based on the following topic:{text} for my powerpoint?"
#     response = generate_response(system_input, question)
    response = modelG.generate_content(f'''System: You are a creative writing assistant. Your mission is to help users generate beautiful and thought provoking powerpoint presentation titles based on the text they input. Generate a title with a very short description of its meaning below. Format like-Title: ... Meaning: ... (The meaning should be written in less than 30 words)
                                    Can you generate a title based on the following topic:{text} for my powerpoint?''')
    return response.text
def generate_table(text): #using neuralchat chatbot
    text=preprocess_textS(text)
    question = f"I'm writing a document on the summary: {text}, and I need a table of contents with only 5 sections to organize the content effectively starting with '1. Introduction' and ending with '5.Conclusion'.. Please generate the table of contents for me with exactly 5 items following the pattern mentioned."
#     response = chatbot.predict(question)
    response = modelG.generate_content(question)
    return response.text
def generate_para(tablec,text): #using neural chat
#     system_input = "You are a creative writing assistant. Your mission is to help users generate detailed information and content based on a given table of contents and input topic. Make it creative,structured with bullet points and paragraphs and detailed information."
#     question = f"Can you generate one elaborate paragraph each for the table of contents {tablec} and based on the reference to the following summary:{text}"
#     response = generate_response(system_input, question)
    response = modelG.generate_content(f'''
    You are a creative writing assistant. Your mission is to help users generate detailed information and content based on a given table of contents and input topic. Make it creative,structured with bullet points and paragraphs and detailed information."
    question = f"Can you generate one elaborate paragraph each for the table of contents {tablec} and based on the reference to the following summary:{text}
    ''')
    return response.text
def generate_paras(tablec,text):  #using gemini api
    response = modelG.generate_content(f"can you generate one big elaborate paragraph of minimum 5 lines for each based on the following table of contents: {tablec} and refering to the summary {text}? Try to include all details present in the summary as well as additional relevant information as you can add. The format should be Slide 1 followed by bullet points, Slide 2 followed by bullet points, etc. for powerpoint presentation. Make it creative with detailed information with minimum 3 points for each slide. YOU MUST NOT make the content in the bullet points bold.")
    return response.text
def get_analysis(sum):
    text = \
    f'''
    I am going to give you a summary report. You have to analyse the sentiment of the summary and then elaborate on it. Make sure you write your analysis in about 100 words.
    Summary: {sum}
    (Start with: In conclusion...)
    '''
    response = modelG.generate_content(text)
    return response.text
def generate_wordcloud(text):
    text = preprocess_textW(text)
    # Generate the word cloud using the default font
    wordcloud = WordCloud(width=800, height=400, background_color='white').generate(text)
    return wordcloud
def generate_keywords(summary, num_keywords=6):
    # Tokenize the summary into words
    words = preprocess_textK(summary)
    # Join the preprocessed words back into a string
    processed_text = " ".join(words)
    # Create TF-IDF vectorizer
    vectorizer = TfidfVectorizer()
    # Fit the vectorizer on the summary
    vectorizer.fit([processed_text])
    # Transform the summary into TF-IDF matrix
    tfidf_matrix = vectorizer.transform([processed_text])
    # Extract feature names (words)
    feature_names = vectorizer.get_feature_names_out()
    # Compute TF-IDF scores for each word
    scores = tfidf_matrix.toarray().flatten()
    # Sort the words based on TF-IDF scores
    keywords = [feature_names[i] for i in scores.argsort()[::-1][:num_keywords]]
    return keywords

def fill_ppt_placeholders(presentation_path, placeholders):
    prs = Presentation(presentation_path)
    current_index = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholder = shape.placeholder_format.idx
                if current_index < len(placeholders):
                    if isinstance(placeholders[current_index], str):
                        if placeholders[current_index].endswith(('.png', '.jpg', '.jpeg', '.gif')):
                            if current_index == len(placeholders) - 1:
                                left_inch = Inches(1)
                                top_inch = Inches(2)
                                width_inch = Inches(8)
                                height_inch = Inches(8)
                                pic = slide.shapes.add_picture(placeholders[current_index], left_inch, top_inch, width_inch, height_inch)
                            else:
                                shape.text = placeholders[current_index]
                        else:
                            shape.text = placeholders[current_index]
                        if placeholder == 0:
                            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        elif placeholder > 0:
                            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    else:
                        shape.text = placeholders[current_index]
                    if placeholder in [0, 2, 4, 6, 8, 10, 12, 14]:
                        shape.text = placeholders[current_index]
                    current_index += 1
                    if current_index == len(placeholders):
                        current_index = 0
    prs.save("Preach.pptx")
        
def process_input(input_data, val='Beige'):
#     input_data = input("Enter the path to the file: ")
    if input_data:
        print("Uploaded successfully.")
    else:
        print("Error: No file provided!")
    # Handle input and convert to transcript
    transcript = handle_input(input_data)
    # wordcloud of summary
    wordcloud = generate_wordcloud(transcript)
    wordcloud.to_file('temp_WC.png')
    # Generate summary from transcript
    summary = generate_summary(transcript)   
    # Generate extended summary from transcript
    count = calculate_percentage_of_words(transcript)
    summary_ext = generate_summary_ext(transcript,count)   
    # keywords of summary
    keywords = generate_keywords(summary)
    # title
    title = generate_title(summary)
    # table of contents using neuralchat chatbot
    tablec=generate_table(summary)
    # paras of summary
    paras = generate_paras(tablec,summary)
    # analysis of summary
    analysis = get_analysis(summary_ext)
    content = 'Summary: '+ summary + '\n\n'+ 'Keywords: ' + str(keywords) + '\n\n' + title + '\n\n' + 'Table of Contents:\n'+ tablec +'\n\n'+ paras +  "\nAnalysis: " + analysis
    with open('content.txt', 'w') as file:
        file.write(content)
    # title_regex = r"Title: \"(.*?)\""
    title_regex = r"Title:\s*(.*)"
    meaning_regex = r"Meaning: (.*?)\n"
    summary_regex = r"Summary: (.*?)\n"
    table_of_contents_regex = r"(\d+\.\s*.*)"
    slide_1_header_regex = r"\*Slide 1: (.*?)\*"
    slide_1_content_regex = r"\**Slide 1: .*?\n(.*?)\n\*\*Slide 2:"
    slide_2_header_regex = r"\*Slide 2: (.*?)\*"
    slide_2_content_regex = r"\**Slide 2: .*?\n(.*?)\n\*\*Slide 3:"
    slide_3_header_regex = r"\*Slide 3: (.*?)\*"
    slide_3_content_regex = r"\**Slide 3: .*?\n(.*?)\n\*\*Slide 4:"
    slide_4_header_regex = r"\*Slide 4: (.*?)\*"
    slide_4_content_regex = r"\**Slide 4: .*?\n(.*?)\n\*\*Slide 5:"
    slide_5_header_regex = r"\*Slide 5: (.*?)\*"
    slide_5_content_regex = r"\**Slide 5: .*?\n(.*?)\nAnalysis:"
    analysis_regex = r"Analysis:\s*(.*)"

    title = re.search(title_regex, content).group(1)
    title = re.sub(r':', r':\n', title)
    meaning = re.search(meaning_regex, content).group(1)
    summary = re.search(summary_regex, content).group(1)
    table_of_contents = re.findall(table_of_contents_regex, content)
    table_of_contents = "\n".join(table_of_contents)
    table_of_contents = re.sub(r'^\d+\.\s*', '', table_of_contents, flags=re.MULTILINE)
    slide_1_header = re.search(slide_1_header_regex, content).group(1)
    slide_1_header = slide_1_header.upper()
    slide_1_content = re.search(slide_1_content_regex, content, re.DOTALL).group(1)
    slide_1_content = re.sub(r'^\-\s*', '', slide_1_content, flags=re.MULTILINE)
    slide_2_header = re.search(slide_2_header_regex, content).group(1)
    slide_2_header = slide_2_header.upper()
    slide_2_content = re.search(slide_2_content_regex, content, re.DOTALL).group(1)
    slide_2_content = re.sub(r'^\-\s*', '', slide_2_content, flags=re.MULTILINE)
    slide_3_header = re.search(slide_3_header_regex, content).group(1)
    slide_3_header = slide_3_header.upper()
    slide_3_content = re.search(slide_3_content_regex, content, re.DOTALL).group(1)
    slide_3_content = re.sub(r'^\-\s*', '', slide_3_content, flags=re.MULTILINE)
    slide_4_header = re.search(slide_4_header_regex, content).group(1)
    slide_4_header = slide_4_header.upper()
    slide_4_content = re.search(slide_4_content_regex, content, re.DOTALL).group(1)
    slide_4_content = re.sub(r'^\-\s*', '', slide_4_content, flags=re.MULTILINE)
    slide_5_header = re.search(slide_5_header_regex, content).group(1)
    slide_5_header = slide_5_header.upper()
    slide_5_content = re.search(slide_5_content_regex, content, re.DOTALL).group(1)
    slide_5_content = re.sub(r'^\-\s*', '', slide_5_content, flags=re.MULTILINE)
    analysis = re.search(analysis_regex, content).group(1)

    arr = [
        title,
        meaning,
        summary,
        table_of_contents,
        slide_1_header,
        slide_1_content,
        slide_2_header,
        slide_2_content,
        slide_3_header,
        slide_3_content,
        slide_4_header,
        slide_4_content,
        slide_5_header,
        slide_5_content,
        analysis,
        'temp_WC.png',]

    if val == 'Beige':
        template = "template_beige.pptx"
        fill_ppt_placeholders(template, arr)
    elif val == 'Green':
        template = "template_green.pptx"
        fill_ppt_placeholders(template, arr)
    elif val == 'Blue':
        template = "template_blue.pptx"
        fill_ppt_placeholders(template, arr)
    elif val == 'Grey':
        template = "template_grey.pptx"
        fill_ppt_placeholders(template, arr)
    elif val == 'Red':
        template = "template_red.pptx"
        fill_ppt_placeholders(template, arr)
    else:
        print("Select proper template")

st.title("Preach: Input to PPT Generator")
# Streamlit app
def main():
    # Input file uploader
    uploaded_file = st.file_uploader("Upload file", type=["mp4", "wav", "txt"])
    option = st.selectbox('Select PPT style',('Beige', 'Green', 'Blue', 'Grey', 'Red'))
    if st.button('Submit'):
        if uploaded_file is not None:
            # Process input and generate PPT
            st.markdown(f"### Take a walk while we generate your PPT...")
            st.markdown(f"#### But be quick! because we are fast 🙃")
            process_input(uploaded_file, option)
            # Display generated PPT
            st.markdown(f"## Generated PowerPoint Presentation...")      
            # Read the content of the generated PPT file
            with open('Preach.pptx', "rb") as f:
                ppt_content = f.read()
                st.download_button(
                label="Click to Download!",
                data=ppt_content,
                file_name="Preach.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            st.error("Failed to generate PowerPoint presentation.")

if __name__ == "__main__":
    main()
